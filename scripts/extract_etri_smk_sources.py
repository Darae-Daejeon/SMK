from __future__ import annotations

import subprocess
import shutil
from pathlib import Path
from typing import Iterable

from pptx import Presentation
from pypdf import PdfReader


ROOT = Path(__file__).resolve().parents[1]
DOWNLOADS = Path.home() / "Downloads"
WORK_DIR = ROOT / "tmp" / "etri-smk"

SOURCE_MAP = [
    {
        "slug": "tech-01",
        "tech_id": "3610-2021-02099",
        "path": DOWNLOADS / "1번기술_003TSL-기술이전심의발표자료-최종.pptx",
        "kind": "pptx",
    },
    {
        "slug": "tech-02",
        "tech_id": "3610-2023-00011",
        "path": DOWNLOADS / "2번기술_PDL-기술이전심의발표자료-v1.pptx",
        "kind": "pptx",
    },
    {
        "slug": "tech-03",
        "tech_id": "5310-2023-00421",
        "path": DOWNLOADS / "3번기술_5.DRL-기술이전심의발표자료-v3.pptx",
        "kind": "pptx",
    },
    {
        "slug": "tech-04",
        "tech_id": "5310-2025-00388",
        "path": DOWNLOADS / "5번기술.pdf",
        "kind": "pdf",
    },
]


def validate_sources(sources: Iterable[dict]) -> None:
    for source in sources:
        path = Path(source["path"])
        if not path.is_file():
            raise FileNotFoundError(f"원본 파일을 찾을 수 없습니다: {path}")


def extract_pptx_text(path: Path) -> list[str]:
    presentation = Presentation(path)
    pages: list[str] = []
    for slide in presentation.slides:
        blocks: list[str] = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                text = shape.text.strip()
                if text:
                    blocks.append(text)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    values = [cell.text.strip() for cell in row.cells]
                    blocks.append(" | ".join(values))
        pages.append("\n".join(blocks))
    return pages


def render_pptx(path: Path, output_dir: Path) -> None:
    import win32com.client

    output_dir.mkdir(parents=True, exist_ok=True)
    app = win32com.client.DispatchEx("PowerPoint.Application")
    presentation = None
    try:
        presentation = app.Presentations.Open(str(path.resolve()), WithWindow=False)
        for index in range(1, presentation.Slides.Count + 1):
            destination = output_dir / f"page-{index:02d}.png"
            presentation.Slides(index).Export(str(destination.resolve()), "PNG", 1920, 1080)
    finally:
        if presentation is not None:
            presentation.Close()
        app.Quit()


def extract_pdf_text(path: Path) -> list[str]:
    reader = PdfReader(path)
    return [(page.extract_text() or "").strip() for page in reader.pages]


def build_pdftoppm_command(executable: Path, source: Path, prefix: Path) -> list[str]:
    arguments = [str(executable), "-png", "-r", "150", str(source), str(prefix)]
    if executable.suffix.lower() in {".cmd", ".bat"}:
        return ["cmd.exe", "/d", "/c", *arguments]
    return arguments


def native_pdftoppm_candidate(executable: Path) -> Path:
    return executable.parents[2] / "native" / "poppler" / "Library" / "bin" / "pdftoppm.exe"


def render_pdf(path: Path, output_dir: Path) -> None:
    output_dir.mkdir(parents=True, exist_ok=True)
    prefix = output_dir / "page"
    executable = shutil.which("pdftoppm")
    if not executable:
        raise FileNotFoundError("pdftoppm 실행 파일을 찾을 수 없습니다")
    executable_path = Path(executable)
    if executable_path.suffix.lower() in {".cmd", ".bat"}:
        native_candidate = native_pdftoppm_candidate(executable_path)
        if native_candidate.is_file():
            executable_path = native_candidate
    subprocess.run(
        build_pdftoppm_command(executable_path, path, prefix),
        check=True,
    )
    for old_path in sorted(output_dir.glob("page-*.png")):
        number = int(old_path.stem.split("-")[-1])
        destination = output_dir / f"page-{number:02d}.png"
        if old_path != destination:
            old_path.rename(destination)


def write_source_note(source: dict, pages: list[str]) -> Path:
    note_dir = WORK_DIR / "source-notes"
    note_dir.mkdir(parents=True, exist_ok=True)
    destination = note_dir / f"{source['slug']}.txt"
    lines = [
        f"원본 파일: {Path(source['path']).name}",
        f"ETRI 기술번호: {source['tech_id']}",
        f"페이지 수: {len(pages)}",
        "",
    ]
    for index, text in enumerate(pages, start=1):
        lines.extend([f"===== PAGE {index} =====", text, ""])
    destination.write_text("\n".join(lines), encoding="utf-8")
    return destination


def extract_source(source: dict) -> tuple[Path, Path]:
    path = Path(source["path"])
    render_dir = WORK_DIR / "rendered" / source["slug"]
    if source["kind"] == "pptx":
        pages = extract_pptx_text(path)
        render_pptx(path, render_dir)
    else:
        pages = extract_pdf_text(path)
        render_pdf(path, render_dir)
    note_path = write_source_note(source, pages)
    return note_path, render_dir


def main() -> None:
    validate_sources(SOURCE_MAP)
    for source in SOURCE_MAP:
        note_path, render_dir = extract_source(source)
        page_count = len(list(render_dir.glob("page-*.png")))
        print(f"{source['slug']}: note={note_path} rendered={page_count}")


if __name__ == "__main__":
    main()
